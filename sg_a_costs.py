# sg_a_costs.py
import streamlit as st
import pandas as pd
import altair as alt
from database import Database
from pptx import Presentation
from pptx.util import Inches
import io

def generate_sg_a_costs_pptx(chart, total_cost):
    """PowerPointファイルを生成"""
    ppt = Presentation()
    slide = ppt.slides.add_slide(ppt.slide_layouts[5])
    slide.shapes.title.text = "販管費データ"
    slide.shapes.add_picture(chart, Inches(1), Inches(1.5), width=Inches(6))
    slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(6), Inches(1)).text = f"総販管費: {total_cost:,} 円"
    pptx_file = "sg_a_costs_report.pptx"
    ppt.save(pptx_file)
    return pptx_file

def sg_a_costs_page():
    st.header("販管費管理")

    # 販管費データ登録フォーム
    with st.form("sg_a_costs_form"):
        category = st.text_input("費目")
        amount = st.number_input("金額（円）", min_value=0, step=1000)
        date = st.date_input("日付")
        submitted = st.form_submit_button("登録")
        if submitted:
            Database.execute_query(
                "INSERT INTO sg_a_costs (category, amount, date) VALUES (?, ?, ?)",
                (category, amount, str(date))
            )
            st.success("販管費データを登録しました！")

    # 販管費データ表示
    data = Database.fetch_data("SELECT category, amount, date FROM sg_a_costs")
    df = pd.DataFrame(data, columns=["費目", "金額", "日付"])
    if not df.empty:
        st.dataframe(df)

        # 総販管費を計算
        total_cost = df["金額"].sum()
        st.metric("総販管費", f"{total_cost:,} 円")

        # グラフ表示
        chart = alt.Chart(df).mark_bar().encode(
            x="費目:N",
            y="金額:Q",
            tooltip=["費目", "金額", "日付"]
        )
        st.altair_chart(chart, use_container_width=True)

        # PowerPointエクスポート
        st.subheader("PowerPointエクスポート")
        if st.button("PowerPointファイルをダウンロード"):
            chart_path = "sg_a_costs_chart.png"
            chart.save(chart_path)
            pptx_file = generate_sg_a_costs_pptx(chart_path, total_cost)
            with open(pptx_file, "rb") as file:
                st.download_button(
                    label="PowerPointをダウンロード",
                    data=file,
                    file_name=pptx_file,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
