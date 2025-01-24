import streamlit as st
import pandas as pd
import math
import os
import matplotlib.pyplot as plt
from database import Database
from pptx import Presentation
from pptx.util import Inches

# チャート画像を保存するためのディレクトリを作成
output_dir = os.path.abspath("output_images")
os.makedirs(output_dir, exist_ok=True)


def save_matplotlib_chart(fig, filename):
    """Matplotlibのグラフを画像として保存"""
    filepath = os.path.join(output_dir, filename)
    fig.savefig(filepath, format="png", dpi=300, bbox_inches="tight")
    plt.close(fig)
    return filepath


def generate_profit_pptx(monthly_chart_path):
    """PowerPointファイルを生成"""
    ppt = Presentation()

    # スライド1: 月毎利益
    slide = ppt.slides.add_slide(ppt.slide_layouts[5])
    slide.shapes.title.text = "月毎利益"
    slide.shapes.add_picture(monthly_chart_path, Inches(1), Inches(1.5), width=Inches(6))

    pptx_file = os.path.join(output_dir, "profit_report.pptx")
    ppt.save(pptx_file)
    return pptx_file


def profit_management_page():
    st.header("利益管理")

    # 売上、原価、販管費データを取得
    sales_total = Database.fetch_data("SELECT SUM(revenue) FROM sales")[0][0] or 0
    cost_total = Database.fetch_data("SELECT SUM(cost) FROM costs")[0][0] or 0
    sg_a_cost_total = Database.fetch_data("SELECT SUM(amount) FROM sg_a_costs")[0][0] or 0

    # 営業利益を計算 (千円単位、小数点切り捨て)
    sales_total = math.floor(sales_total / 1000)
    cost_total = math.floor(cost_total / 1000)
    sg_a_cost_total = math.floor(sg_a_cost_total / 1000)
    total_profit = sales_total - cost_total - sg_a_cost_total

    # 利益情報を表示
    st.metric("総売上", f"{sales_total:,} 千円")
    st.metric("総原価", f"{cost_total:,} 千円")
    st.metric("総販管費", f"{sg_a_cost_total:,} 千円")
    st.metric("営業利益", f"{total_profit:,} 千円")

    # 仮の月毎利益データ (ここをデータベースのクエリに変更可能)
    monthly_data = pd.DataFrame({
        "月": ["2023-01", "2023-02", "2023-03"],
        "利益": [total_profit, total_profit * 0.8, total_profit * 1.1]
    })

    # 月毎利益グラフ（Matplotlib）
    fig_monthly, ax_monthly = plt.subplots(figsize=(10, 6))

    # 折れ線グラフを描画
    line, = ax_monthly.plot(
        monthly_data["月"],
        monthly_data["利益"],
        marker="o",
        linestyle="-",
        color="blue"
    )

    # 凡例を手動で設定
    ax_monthly.legend([line], ["営業利益"], fontsize=12, loc="upper left", frameon=True, shadow=True)

    # グラフのタイトル、ラベル、凡例を設定
    ax_monthly.set_title(
        "月毎利益（単位：千円）",
        fontsize=20, weight="bold", color="darkblue", pad=20
    )
    ax_monthly.set_ylabel(
        "利益（千円）",
        fontsize=14, labelpad=10, fontweight="bold", color="black"
    )
    ax_monthly.set_xlabel(
        "月次推移",
        fontsize=14, labelpad=10, fontweight="bold", color="black"
    )
    ax_monthly.tick_params(axis="both", which="major", labelsize=12)
    ax_monthly.grid(True, linestyle="--", alpha=0.5)

    # グラフ画像を保存
    monthly_chart_path = save_matplotlib_chart(fig_monthly, "monthly_chart.png")

    # Streamlitにグラフを表示
    st.image(monthly_chart_path, caption="月毎利益（単位：千円）", use_container_width=True)

    # PowerPointエクスポート
    st.subheader("PowerPointエクスポート")
    if st.button("PowerPointファイルをダウンロード"):
        pptx_file = generate_profit_pptx(monthly_chart_path)
        with open(pptx_file, "rb") as file:
            st.download_button(
                label="PowerPointをダウンロード",
                data=file,
                file_name="profit_report.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
