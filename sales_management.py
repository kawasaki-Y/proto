import streamlit as st
import pandas as pd
import math
import os
import matplotlib.pyplot as plt
from database import Database
from pptx import Presentation
from pptx.util import Inches

# 出力ディレクトリを設定
output_dir = os.path.abspath("output_images")
os.makedirs(output_dir, exist_ok=True)

def save_matplotlib_chart(fig, filename):
    """Matplotlibのグラフを画像として保存"""
    filepath = os.path.join(output_dir, filename)
    fig.savefig(filepath, format="png", dpi=300, bbox_inches="tight")
    plt.close(fig)
    return filepath

def generate_sales_pptx(bar_chart_path, line_chart_path, pie_chart_path, total_sales, sales_difference):
    """PowerPointファイルを生成"""
    ppt = Presentation()

    # スライド1: 棒グラフ
    slide = ppt.slides.add_slide(ppt.slide_layouts[5])
    slide.shapes.title.text = "目標売上と総売上の比較"
    slide.shapes.add_picture(bar_chart_path, Inches(1), Inches(1.5), width=Inches(6))

    # スライド2: 折れ線グラフ
    slide = ppt.slides.add_slide(ppt.slide_layouts[5])
    slide.shapes.title.text = "月別総売上"
    slide.shapes.add_picture(line_chart_path, Inches(1), Inches(1.5), width=Inches(6))

    # スライド3: 円グラフ
    slide = ppt.slides.add_slide(ppt.slide_layouts[5])
    slide.shapes.title.text = "タグごとの売上割合"
    slide.shapes.add_picture(pie_chart_path, Inches(1), Inches(1.5), width=Inches(6))

    # スライド4: 総売上と目標差分
    slide = ppt.slides.add_slide(ppt.slide_layouts[5])
    slide.shapes.title.text = "総売上と目標差分"
    textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    textbox.text = f"総売上: {total_sales:,} 千円\n目標との差分: {sales_difference:+,} 千円"

    pptx_file = os.path.join(output_dir, "sales_report.pptx")
    ppt.save(pptx_file)
    return pptx_file

def sales_management_page():
    st.header("売上管理")

    # 登録済みタグを取得
    tags_data = Database.fetch_data("SELECT tag_name FROM tags")
    tags = [tag[0] for tag in tags_data]

    # 売上データ登録フォーム
    with st.form("sales_form"):
        project = st.text_input("案件名")
        tag = st.selectbox("タグ", options=tags if tags else ["タグなし"])
        revenue = st.number_input("売上金額（千円単位）", min_value=0, step=1)  # 千円単位で入力
        date = st.date_input("日付")
        if st.form_submit_button("登録"):
            Database.execute_query(
                "INSERT INTO sales (project, tag, revenue, date) VALUES (?, ?, ?, ?)",
                (project, tag, revenue, str(date))
            )
            st.success("売上データを登録しました！")
            st.experimental_rerun()

    # 売上データ取得
    data = Database.fetch_data("SELECT id, project, tag, revenue, date FROM sales")
    df = pd.DataFrame(data, columns=["ID", "案件名", "タグ", "売上金額", "日付"])
    df["売上金額"] = df["売上金額"].apply(math.floor)
    df.insert(0, "番号", range(1, len(df) + 1))

    if not df.empty:
        st.subheader("売上データ一覧")
        st.dataframe(df[["番号", "案件名", "タグ", "売上金額", "日付"]], use_container_width=True)

        # 目標売上の取得
        target_revenue_data = Database.fetch_data("SELECT amount FROM target_revenue")
        target_revenue = math.floor(target_revenue_data[0][0]) if target_revenue_data else 0

        # 総売上の計算
        total_sales = df["売上金額"].sum()
        sales_difference = total_sales - target_revenue

        # メトリクス表示
        st.metric("目標売上", f"{target_revenue:,} 千円")
        st.metric("総売上", f"{total_sales:,} 千円")
        st.metric("目標との差分", f"{sales_difference:+,} 千円")

        # データ編集と削除
        st.subheader("データ編集・削除")
        selected_row = st.selectbox(
            "編集・削除するデータを選択",
            options=df.to_dict(orient="records"),
            format_func=lambda x: f"{x['案件名']} - {x['売上金額']}千円 - {x['日付']}"
        )
        if selected_row:
            edit_project = st.text_input("新しい案件名", value=selected_row["案件名"], key="edit_project")
            edit_revenue = st.number_input("新しい売上金額（千円単位）", value=selected_row["売上金額"], min_value=0, step=1, key="edit_revenue")
            edit_date = st.date_input("新しい日付", value=pd.to_datetime(selected_row["日付"]), key="edit_date")

            if st.button("変更を保存"):
                Database.execute_query(
                    "UPDATE sales SET project = ?, revenue = ?, date = ? WHERE id = ?",
                    (edit_project, edit_revenue, str(edit_date), selected_row["ID"])
                )
                st.success("データを更新しました！")
                st.experimental_rerun()

            if st.button("データを削除"):
                Database.execute_query("DELETE FROM sales WHERE id = ?", (selected_row["ID"],))
                st.success("データを削除しました！")
                st.experimental_rerun()

        # グラフ作成
        st.subheader("売上データのグラフ")
        df["月"] = pd.to_datetime(df["日付"]).dt.strftime("%Y-%m")

        # 棒グラフ
        fig_bar, ax_bar = plt.subplots(figsize=(8, 6))
        comparison_data = {"項目": ["目標売上", "総売上"], "金額": [target_revenue, total_sales]}
        pd.DataFrame(comparison_data).set_index("項目")["金額"].plot(
            kind="bar", ax=ax_bar, color=["#FF7F50", "#4682B4"]
        )
        ax_bar.set_title("目標売上と総売上の比較", fontsize=20, fontweight="bold")
        ax_bar.set_xlabel("項目", fontsize=14)
        ax_bar.set_ylabel("売上高（千円）", fontsize=14)
        for i, v in enumerate([target_revenue, total_sales]):
            ax_bar.text(i, v + 500, f"{v:,} 千円", ha="center", fontsize=12)
        bar_chart_path = save_matplotlib_chart(fig_bar, "bar_chart.png")
        st.image(bar_chart_path, caption="目標売上と総売上の比較", use_container_width=True)

        # 折れ線グラフ
        fig_line, ax_line = plt.subplots(figsize=(8, 6))
        monthly_sales = df.groupby("月")["売上金額"].sum()
        ax_line.plot(monthly_sales.index, monthly_sales.values, marker="o", linestyle="-", color="#2E8B57")
        ax_line.set_title("月別総売上", fontsize=20, fontweight="bold")
        ax_line.set_xlabel("月次推移", fontsize=14)
        ax_line.set_ylabel("売上高（千円）", fontsize=14)
        for x, y in zip(monthly_sales.index, monthly_sales.values):
            ax_line.text(x, y + 500, f"{y:,} 千円", fontsize=10, ha="center")
        line_chart_path = save_matplotlib_chart(fig_line, "line_chart.png")
        st.image(line_chart_path, caption="月別総売上", use_container_width=True)

        # 円グラフ
        fig_pie, ax_pie = plt.subplots(figsize=(6, 6))
        tag_sales = df.groupby("タグ")["売上金額"].sum()
        tag_sales.plot(
            kind="pie", ax=ax_pie, autopct="%1.1f%%", startangle=90, colors=["#FF7F50", "#4682B4", "#32CD32"]
        )
        ax_pie.set_title("タグごとの売上割合", fontsize=20, fontweight="bold")
        ax_pie.set_ylabel("")
        pie_chart_path = save_matplotlib_chart(fig_pie, "pie_chart.png")
        st.image(pie_chart_path, caption="タグごとの売上割合", use_container_width=True)

        # PowerPointエクスポート
        st.subheader("PowerPointエクスポート")
        if st.button("PowerPointファイルをダウンロード"):
            pptx_file = generate_sales_pptx(bar_chart_path, line_chart_path, pie_chart_path, total_sales, sales_difference)
            with open(pptx_file, "rb") as file:
                st.download_button(
                    label="PowerPointをダウンロード",
                    data=file,
                    file_name="sales_report.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
