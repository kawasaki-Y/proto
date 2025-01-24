# cost_management.py
import streamlit as st
import pandas as pd
import altair as alt
from database import Database
from pptx import Presentation
from pptx.util import Inches
import io

def generate_cost_pptx(chart, total_cost):
    """PowerPointファイルを生成"""
    ppt = Presentation()
    slide = ppt.slides.add_slide(ppt.slide_layouts[5])
    slide.shapes.title.text = "原価データ"
    slide.shapes.add_picture(chart, Inches(1), Inches(1.5), width=Inches(6))
    slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(6), Inches(1)).text = f"総原価: {total_cost:,} 円"
    pptx_file = "cost_report.pptx"
    ppt.save(pptx_file)
    return pptx_file

def cost_management_page():
    st.header("原価管理")

    # 原価データ登録フォーム
    with st.form("cost_form"):
        project = st.text_input("案件名")
        cost = st.number_input("原価金額（円）", min_value=0, step=1000)
        date = st.date_input("日付")
        submitted = st.form_submit_button("登録")
        if submitted:
            Database.execute_query(
                "INSERT INTO costs (project, cost, date) VALUES (?, ?, ?)",
                (project, cost, str(date))
            )
            st.success("原価データを登録しました！")

    # 原価データ表示
    data = Database.fetch_data("SELECT project, cost, date FROM costs")
    df = pd.DataFrame(data, columns=["案件名", "原価金額", "日付"])
    if not df.empty:
        st.dataframe(df)

        # 総原価を計算
        total_cost = df["原価金額"].sum()
        st.metric("総原価", f"{total_cost:,} 円")

        # グラフ表示
        chart = alt.Chart(df).mark_bar().encode(
            x="案件名:N",
            y="原価金額:Q",
            tooltip=["案件名", "原価金額", "日付"]
        )
        st.altair_chart(chart, use_container_width=True)

        # PowerPointエクスポート
        st.subheader("PowerPointエクスポート")
        if st.button("PowerPointファイルをダウンロード"):
            chart_path = "cost_chart.png"
            chart.save(chart_path)
            pptx_file = generate_cost_pptx(chart_path, total_cost)
            with open(pptx_file, "rb") as file:
                st.download_button(
                    label="PowerPointをダウンロード",
                    data=file,
                    file_name=pptx_file,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
